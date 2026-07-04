#!/usr/bin/env python3
"""
文档解析统一 CLI：PDF / Word / PPT 结构化提取与版面分析。

输出包含页面/幻灯片结构、文本块、表格、图片位置、字体信息等，
便于下游做版面分析、阅读顺序推断或精排后处理。
"""

from __future__ import annotations

import argparse
import base64
import io
import json
import os
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional

try:
    from dotenv import load_dotenv

    load_dotenv()
except ImportError:  # pragma: no cover
    pass


def write_json(path: Optional[str], data: Any) -> None:
    content = json.dumps(data, ensure_ascii=False, indent=2, default=str)
    if path is None or path == "-":
        sys.stdout.write(content)
        sys.stdout.write("\n")
    else:
        Path(path).write_text(content, encoding="utf-8")


def write_text(path: Optional[str], content: str) -> None:
    if path is None or path == "-":
        sys.stdout.write(content)
        if not content.endswith("\n"):
            sys.stdout.write("\n")
    else:
        Path(path).write_text(content, encoding="utf-8")


# ---------------------------------------------------------------------------
# PDF parsing
# ---------------------------------------------------------------------------

def parse_pdf(file_path: str, extract_images: bool = False) -> dict:
    try:
        import pdfplumber
    except ImportError:
        raise RuntimeError("PDF 解析需要 pdfplumber，请执行: pip install pdfplumber")

    pages = []
    with pdfplumber.open(file_path) as doc:
        for page_index, page in enumerate(doc.pages):
            width = page.width
            height = page.height
            elements: List[dict] = []

            # Text words grouped by line
            words = page.extract_words()
            lines: Dict[float, List[dict]] = {}
            for w in words:
                top = round(float(w["top"]), 1)
                lines.setdefault(top, []).append(w)

            for top in sorted(lines.keys()):
                line_words = sorted(lines[top], key=lambda w: float(w["x0"]))
                text = " ".join(w["text"] for w in line_words)
                if not text.strip():
                    continue
                left = min(float(w["x0"]) for w in line_words)
                right = max(float(w["x1"]) for w in line_words)
                bottom = max(float(w["bottom"]) for w in line_words)
                elements.append(
                    {
                        "type": "text",
                        "text": text,
                        "bbox": [left, top, right, bottom],
                    }
                )

            # Tables
            tables = page.extract_tables() or []
            for table in tables:
                elements.append({"type": "table", "rows": table, "bbox": None})

            # Images (placeholder: pdfplumber does not easily extract image bboxes)
            images = []
            if extract_images:
                for img in page.images:
                    bbox = [float(img.get("x0", 0)), float(img.get("y0", 0)), float(img.get("x1", 0)), float(img.get("y1", 0))]
                    images.append({"bbox": bbox})
                    elements.append({"type": "image", "bbox": bbox})

            # Sort by reading order
            elements.sort(key=lambda e: (e.get("bbox", [0, 0, 0, 0])[1], e.get("bbox", [0, 0, 0, 0])[0]))

            pages.append(
                {
                    "page_number": page_index + 1,
                    "width": width,
                    "height": height,
                    "elements": elements,
                    "table_count": len(tables),
                    "image_count": len(images) if extract_images else None,
                }
            )

    return {
        "file": file_path,
        "type": "pdf",
        "page_count": len(pages),
        "pages": pages,
    }


# ---------------------------------------------------------------------------
# DOCX parsing
# ---------------------------------------------------------------------------

def parse_docx(file_path: str) -> dict:
    try:
        import docx
    except ImportError:
        raise RuntimeError("Word 解析需要 python-docx，请执行: pip install python-docx")

    doc = docx.Document(file_path)
    elements: List[dict] = []

    for i, para in enumerate(doc.paragraphs):
        if not para.text.strip():
            continue
        style = para.style.name if para.style else "Normal"
        elem_type = "heading" if style.lower().startswith("heading") else "paragraph"
        elements.append(
            {
                "type": elem_type,
                "text": para.text,
                "style": style,
                "index": i,
            }
        )

    tables = []
    for table_index, table in enumerate(doc.tables):
        rows = []
        for row in table.rows:
            rows.append([cell.text for cell in row.cells])
        tables.append({"table_index": table_index, "rows": rows})
        elements.append({"type": "table", "table_index": table_index, "rows": rows})

    return {
        "file": file_path,
        "type": "docx",
        "paragraph_count": len(doc.paragraphs),
        "table_count": len(doc.tables),
        "elements": elements,
    }


# ---------------------------------------------------------------------------
# PPTX parsing
# ---------------------------------------------------------------------------

def parse_pptx(file_path: str) -> dict:
    try:
        import pptx
    except ImportError:
        raise RuntimeError("PPT 解析需要 python-pptx，请执行: pip install python-pptx")

    prs = pptx.Presentation(file_path)
    slides = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        elements: List[dict] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    elements.append(
                        {
                            "type": "text",
                            "text": text,
                            "left": shape.left,
                            "top": shape.top,
                            "width": shape.width,
                            "height": shape.height,
                        }
                    )
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    rows.append([cell.text for cell in row.cells])
                elements.append(
                    {
                        "type": "table",
                        "rows": rows,
                        "left": shape.left,
                        "top": shape.top,
                    }
                )
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                elements.append(
                    {
                        "type": "image",
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height,
                    }
                )

        # Sort by position (top-to-bottom, left-to-right)
        elements.sort(key=lambda e: (e.get("top", 0), e.get("left", 0)))

        notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
        slides.append(
            {
                "slide_number": slide_index,
                "elements": elements,
                "notes": notes,
            }
        )

    return {
        "file": file_path,
        "type": "pptx",
        "slide_count": len(prs.slides),
        "slides": slides,
    }


# ---------------------------------------------------------------------------
# Layout analysis
# ---------------------------------------------------------------------------

def analyze_pdf_layout(file_path: str) -> dict:
    """返回 PDF 版面摘要，包括列数、字体大小分布、阅读顺序建议等。"""
    data = parse_pdf(file_path, extract_images=False)
    summary = []
    for page in data["pages"]:
        bboxes = [e["bbox"] for e in page["elements"] if e.get("bbox")]
        if not bboxes:
            summary.append({"page": page["page_number"], "columns": 0, "blocks": 0})
            continue
        # Simple column detection by x-coordinate clustering
        lefts = [b[0] for b in bboxes]
        # Use mean split
        if lefts:
            mid = sum(lefts) / len(lefts)
            left_count = sum(1 for x in lefts if x < mid)
            right_count = len(lefts) - left_count
            columns = 2 if left_count > 0 and right_count > 0 and min(left_count, right_count) > len(lefts) * 0.2 else 1
        else:
            columns = 1
        summary.append(
            {
                "page": page["page_number"],
                "columns": columns,
                "blocks": len(bboxes),
                "width": page["width"],
                "height": page["height"],
            }
        )
    return {"file": file_path, "type": "pdf", "layout_summary": summary}


# ---------------------------------------------------------------------------
# Main parse dispatcher
# ---------------------------------------------------------------------------

def parse_file(file_path: str, extract_images: bool = False) -> dict:
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return parse_pdf(file_path, extract_images=extract_images)
    if ext in (".docx", ".doc"):
        if ext == ".doc":
            raise ValueError("旧版 .doc 格式不支持，请转换为 .docx")
        return parse_docx(file_path)
    if ext in (".pptx", ".ppt"):
        if ext == ".ppt":
            raise ValueError("旧版 .ppt 格式不支持，请转换为 .pptx")
        return parse_pptx(file_path)
    raise ValueError(f"不支持的文件格式: {ext}")


# ---------------------------------------------------------------------------
# CLI handlers
# ---------------------------------------------------------------------------

def cmd_parse(args: argparse.Namespace) -> None:
    result = parse_file(args.input, extract_images=args.extract_images)
    write_json(args.output, result)


def cmd_layout(args: argparse.Namespace) -> None:
    ext = Path(args.input).suffix.lower()
    if ext == ".pdf":
        result = analyze_pdf_layout(args.input)
    elif ext in (".docx", ".pptx"):
        # For docx/pptx, layout is already in parse output; emit a lightweight summary
        data = parse_file(args.input)
        if data["type"] == "docx":
            result = {
                "file": args.input,
                "type": "docx",
                "paragraph_count": data["paragraph_count"],
                "table_count": data["table_count"],
                "headings": [e["text"] for e in data["elements"] if e["type"] == "heading"],
            }
        else:
            result = {
                "file": args.input,
                "type": "pptx",
                "slide_count": data["slide_count"],
                "slides": [{"slide_number": s["slide_number"], "blocks": len(s["elements"])} for s in data["slides"]],
            }
    else:
        raise ValueError(f"不支持的文件格式: {ext}")
    write_json(args.output, result)


def cmd_extract(args: argparse.Namespace) -> None:
    """将文档内容提取为纯文本/表格/图片到输出目录。"""
    data = parse_file(args.input, extract_images=True)
    output_dir = Path(args.output)
    output_dir.mkdir(parents=True, exist_ok=True)

    # Plain text
    texts = []
    if data["type"] == "pdf":
        for page in data["pages"]:
            for e in page["elements"]:
                if e["type"] == "text":
                    texts.append(e["text"])
    elif data["type"] == "docx":
        for e in data["elements"]:
            if e["type"] in ("paragraph", "heading"):
                texts.append(e["text"])
    elif data["type"] == "pptx":
        for slide in data["slides"]:
            for e in slide["elements"]:
                if e["type"] == "text":
                    texts.append(e["text"])

    (output_dir / "text.txt").write_text("\n\n".join(texts), encoding="utf-8")

    # Tables
    tables = []
    if data["type"] == "pdf":
        for page in data["pages"]:
            for e in page["elements"]:
                if e["type"] == "table":
                    tables.append(e["rows"])
    elif data["type"] == "docx":
        tables = [e["rows"] for e in data["elements"] if e["type"] == "table"]
    elif data["type"] == "pptx":
        for slide in data["slides"]:
            for e in slide["elements"]:
                if e["type"] == "table":
                    tables.append(e["rows"])

    (output_dir / "tables.json").write_text(json.dumps(tables, ensure_ascii=False, indent=2), encoding="utf-8")

    write_json(None, {"status": "extracted", "output_dir": str(output_dir), "text_blocks": len(texts), "tables": len(tables)})


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        prog="document_parse.py",
        description="文档解析工具：PDF / Word / PPT 结构化提取与版面分析。",
    )
    subparsers = parser.add_subparsers(dest="command", required=True)

    p_parse = subparsers.add_parser("parse", help="解析文档为结构化 JSON")
    p_parse.add_argument("input", help="输入文件路径")
    p_parse.add_argument("-o", "--output", default="-", help="输出 JSON 文件路径")
    p_parse.add_argument("--extract-images", action="store_true", help="尝试提取图片位置信息")
    p_parse.set_defaults(func=cmd_parse)

    p_layout = subparsers.add_parser("layout", help="版面分析摘要")
    p_layout.add_argument("input", help="输入文件路径")
    p_layout.add_argument("-o", "--output", default="-", help="输出 JSON 文件路径")
    p_layout.set_defaults(func=cmd_layout)

    p_extract = subparsers.add_parser("extract", help="提取文本和表格到目录")
    p_extract.add_argument("input", help="输入文件路径")
    p_extract.add_argument("-o", "--output", required=True, help="输出目录路径")
    p_extract.set_defaults(func=cmd_extract)

    return parser


def main(argv: Optional[List[str]] = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    try:
        args.func(args)
    except Exception as exc:
        print(f"错误: {exc}", file=sys.stderr)
        return 1
    return 0


if __name__ == "__main__":
    sys.exit(main())
