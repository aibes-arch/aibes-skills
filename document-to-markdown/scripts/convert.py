#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Convert common office documents to Markdown.

Supported formats:
  - Word: .docx (via python-docx)
  - PDF: .pdf (via PyMuPDF)
  - PowerPoint: .pptx (via python-pptx)
  - Excel: .xlsx (via openpyxl)
  - Fallback: anything pandoc can read (requires pandoc installed)
"""

import argparse
import re
import sys
from pathlib import Path

try:
    from docx import Document
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    HAS_DOCX = True
except Exception:  # pragma: no cover
    Document = qn = Table = Paragraph = None
    HAS_DOCX = False

try:
    import fitz  # PyMuPDF

    HAS_FITZ = True
except Exception:  # pragma: no cover
    fitz = None
    HAS_FITZ = False

try:
    from pptx import Presentation

    HAS_PPTX = True
except Exception:  # pragma: no cover
    Presentation = None
    HAS_PPTX = False

try:
    from openpyxl import load_workbook

    HAS_XLSX = True
except Exception:  # pragma: no cover
    load_workbook = None
    HAS_XLSX = False

try:
    import pypandoc

    HAS_PYPANDOC = True
except Exception:  # pragma: no cover
    pypandoc = None
    HAS_PYPANDOC = False


SUPPORTED_EXTS = {".docx", ".pdf", ".pptx", ".xlsx"}
# Additional formats handled via pandoc fallback when available.
PANDOC_INPUT_EXTS = {
    ".doc",
    ".odt",
    ".rtf",
    ".html",
    ".htm",
    ".tex",
    ".epub",
    ".txt",
    ".md",
}
ALL_SUPPORTED_EXTS = SUPPORTED_EXTS | PANDOC_INPUT_EXTS


def sanitize_filename(name: str) -> str:
    """Remove characters that are illegal in Windows filenames."""
    name = re.sub(r'[\\/*?:"<>|]', "_", name)
    return name.strip(" .") or "output"


def write_output(text: str, output_path: Path) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(text, encoding="utf-8")
    return output_path


def _escape_md(text: str) -> str:
    """Escape Markdown special characters inside a table cell."""
    return text.replace("|", "\\|").replace("\n", "<br>")


def _convert_run(run) -> str:
    text = run.text
    if not text:
        return ""
    if run.bold and run.italic:
        return f"***{text}***"
    if run.bold:
        return f"**{text}**"
    if run.italic:
        return f"*{text}*"
    return text


def _convert_paragraph(paragraph, image_map: dict) -> str:
    style_name = paragraph.style.name

    # Heading
    if style_name.startswith("Heading "):
        level = 1
        try:
            level = int(style_name.split()[-1])
        except ValueError:
            pass
        return "#" * level + " " + "".join(_convert_run(r) for r in paragraph.runs)

    # Lists
    if "List Bullet" in style_name:
        return "- " + "".join(_convert_run(r) for r in paragraph.runs)
    if "List Number" in style_name:
        return "1. " + "".join(_convert_run(r) for r in paragraph.runs)

    # Inline images and text
    pieces = []
    for run in paragraph.runs:
        blips = run._r.xpath(".//a:blip")
        if blips:
            embed = blips[0].get(qn("r:embed"))
            if embed and embed in image_map:
                pieces.append(f"![image]({image_map[embed]})")
        pieces.append(_convert_run(run))
    return "".join(pieces)


def _convert_table(table) -> str:
    rows = []
    for row in table.rows:
        cells = [_escape_md(cell.text.strip()) for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
    if rows:
        col_count = len(table.rows[0].cells)
        rows.insert(1, "| " + " | ".join(["---"] * col_count) + " |")
    return "\n".join(rows)


def convert_docx(input_path: Path, output_dir: Path) -> str:
    """Convert a .docx file to Markdown, extracting images to images/."""
    doc = Document(input_path)
    images_dir = output_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    # Extract embedded images and build rId -> relative path map
    image_map = {}
    image_counter = 0
    for rel in doc.part.rels.values():
        if "image" not in rel.reltype:
            continue
        image_counter += 1
        ext = Path(rel.target_ref).suffix.lstrip(".").lower()
        if ext in ("jpeg", "jpg"):
            ext = "jpg"
        if not ext:
            ext = "png"
        filename = f"image_{image_counter:04d}.{ext}"
        image_path = images_dir / filename
        with open(image_path, "wb") as f:
            f.write(rel.target_part.blob)
        image_map[rel.rId] = f"./images/{filename}"

    parts = []
    for child in doc.element.body:
        if child.tag == qn("w:p"):
            paragraph = Paragraph(child, doc)
            text = _convert_paragraph(paragraph, image_map)
            if text.strip():
                parts.append(text)
        elif child.tag == qn("w:tbl"):
            table = Table(child, doc)
            parts.append(_convert_table(table))

    # Compact adjacent list items into a single Markdown list block.
    compacted = []
    for part in parts:
        if compacted and (part.startswith(("- ", "1. "))) and (
            compacted[-1].startswith(("- ", "1. "))
        ):
            compacted[-1] += "\n" + part
        else:
            compacted.append(part)

    return "\n\n".join(compacted)


def convert_pdf(input_path: Path, page_separator: str) -> str:
    """Convert a .pdf file to Markdown using PyMuPDF."""
    doc = fitz.open(input_path)
    pages = []
    for page in doc:
        # PyMuPDF 1.23+ supports native markdown output.
        try:
            text = page.get_text("markdown")
        except Exception:
            text = ""
        if not text or not text.strip():
            text = page.get_text()
        pages.append(text.strip())
    return page_separator.join(pages)


def convert_pptx(input_path: Path, page_separator: str) -> str:
    """Convert a .pptx file to Markdown, one slide per section."""
    prs = Presentation(input_path)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            txt = shape.text.strip()
            if txt:
                texts.append(txt)
        if texts:
            slides.append(f"## Slide {i}\n\n" + "\n\n".join(texts))
    return page_separator.join(slides)


def convert_xlsx(input_path: Path) -> str:
    """Convert a .xlsx file to Markdown tables, one per worksheet."""
    wb = load_workbook(input_path, data_only=True)
    sheets = []
    for sheet in wb.worksheets:
        rows = []
        max_col = 0
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell) if cell is not None else "" for cell in row]
            max_col = max(max_col, len(cells))
            rows.append("| " + " | ".join(cells) + " |")
        if rows:
            header_sep = "| " + " | ".join(["---"] * max_col) + " |"
            rows.insert(1, header_sep)
            sheets.append(f"## Sheet: {sheet.title}\n\n" + "\n".join(rows))
    return "\n\n".join(sheets)


def convert_with_pandoc(input_path: Path, to: str = "markdown") -> str:
    """Fallback conversion via pandoc."""
    return pypandoc.convert_file(str(input_path), to, format=None)


def convert_file(input_path: Path, output_path: Path, args) -> Path:
    """Dispatch to the appropriate converter based on file extension."""
    ext = input_path.suffix.lower()
    output_dir = output_path.parent

    if ext == ".docx" and HAS_DOCX:
        text = convert_docx(input_path, output_dir)
    elif ext == ".pdf" and HAS_FITZ:
        text = convert_pdf(input_path, args.page_separator)
    elif ext == ".pptx" and HAS_PPTX:
        text = convert_pptx(input_path, args.page_separator)
    elif ext == ".xlsx" and HAS_XLSX:
        text = convert_xlsx(input_path)
    elif HAS_PYPANDOC:
        try:
            text = convert_with_pandoc(input_path)
        except Exception as exc:
            raise RuntimeError(f"Pandoc fallback failed for {input_path}: {exc}")
    else:
        raise RuntimeError(f"Unsupported file format: {ext}")

    return write_output(text, output_path)


def convert_directory(input_dir: Path, output_dir: Path, args):
    """Convert every supported file in a directory."""
    files = [p for p in input_dir.iterdir() if p.is_file()]
    results = []
    for file in files:
        ext = file.suffix.lower()
        if ext not in ALL_SUPPORTED_EXTS:
            continue
        if ext not in SUPPORTED_EXTS and not HAS_PYPANDOC:
            continue
        base = sanitize_filename(file.stem)
        md_name = base + ".md"
        out = output_dir / md_name
        if out.exists():
            for n in range(1, 1000):
                candidate = output_dir / f"{base}_{n:03d}.md"
                if not candidate.exists():
                    out = candidate
                    break
        try:
            results.append(convert_file(file, out, args))
        except Exception as exc:
            print(f"[ERROR] {file}: {exc}", file=sys.stderr)
    return results


def main() -> int:
    parser = argparse.ArgumentParser(description="Convert documents to Markdown")
    parser.add_argument("input", help="Input file or directory")
    parser.add_argument("-o", "--output", help="Output Markdown file or directory")
    parser.add_argument(
        "--page-separator",
        default="\n\n---\n\n",
        help="Separator between PDF/PPT pages or slides",
    )
    args = parser.parse_args()

    input_path = Path(args.input).resolve()
    if not input_path.exists():
        print(f"Error: {input_path} does not exist", file=sys.stderr)
        return 1

    if input_path.is_file():
        output_path = (
            Path(args.output).resolve()
            if args.output
            else input_path.with_suffix(".md")
        )
        try:
            convert_file(input_path, output_path, args)
            print(output_path)
        except Exception as exc:
            print(f"Error: {exc}", file=sys.stderr)
            return 1
    else:
        output_dir = Path(args.output).resolve() if args.output else input_path
        output_dir.mkdir(parents=True, exist_ok=True)
        for result in convert_directory(input_path, output_dir, args):
            print(result)

    return 0


if __name__ == "__main__":
    sys.exit(main())
